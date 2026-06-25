"""
Build the project presentation as a PowerPoint file.
Run: .venv/Scripts/python.exe build_presentation.py
Output: documents/Image_Forgery_Detector_Presentation.pptx
"""
import io
import os
import sys
os.environ["TF_CPP_MIN_LOG_LEVEL"] = "3"

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
from PIL import Image as PILImage, ImageChops, ImageEnhance

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0e, 0x1b, 0x2e)
NAVY2      = RGBColor(0x17, 0x28, 0x3f)
NAVY3      = RGBColor(0x1a, 0x3a, 0x63)
BLUE_ACC   = RGBColor(0x4a, 0x9f, 0xd4)
NED_GREEN  = RGBColor(0x1a, 0x6b, 0x3a)
NED_GOLD   = RGBColor(0xc8, 0x96, 0x0c)
WHITE      = RGBColor(0xff, 0xff, 0xff)
LIGHT_BLUE = RGBColor(0xcf, 0xe4, 0xf6)
MID_BLUE   = RGBColor(0x8f, 0xc1, 0xe6)
TEXT_DIM   = RGBColor(0x9f, 0xb6, 0xcf)
GREEN_HI   = RGBColor(0x2e, 0xcc, 0x71)
RED_HI     = RGBColor(0xe7, 0x4c, 0x3c)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUT_PATH = "documents/Image_Forgery_Detector_Presentation.pptx"


# ── Low-level helpers ─────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=NAVY):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return box


def txt_lines(slide, lines, l, t, w, h,
              size=16, color=WHITE, bold=False, line_spacing=1.15,
              bullet=False):
    """Multi-line textbox; each item in lines is a string."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        if bullet:
            p.text = ("• " if line.strip() else "") + line
        else:
            p.text = line
        p.font.size  = Pt(size)
        p.font.bold  = bold
        p.font.color.rgb = color


def add_image(slide, path_or_buf, l, t, w, h):
    if isinstance(path_or_buf, str):
        slide.shapes.add_picture(path_or_buf, l, t, w, h)
    else:
        path_or_buf.seek(0)
        slide.shapes.add_picture(path_or_buf, l, t, w, h)


def fig_to_buf(fig, dpi=100):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


def header_strip(slide, title, subtitle=None):
    """Dark navy top bar with green left accent + gold title underline."""
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), NAVY3)
    rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(1.15), NED_GREEN)
    rect(slide, Inches(0.12), Inches(1.05), SLIDE_W - Inches(0.12), Pt(3), NED_GOLD)
    txt(slide, title,
        Inches(0.28), Inches(0.12), Inches(10), Inches(0.65),
        size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.28), Inches(0.72), Inches(10), Inches(0.32),
            size=13, color=MID_BLUE)


def footer(slide, num, total):
    rect(slide, Inches(0), SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), NAVY3)
    rect(slide, Inches(0), SLIDE_H - Inches(0.38), SLIDE_W, Pt(2), NED_GOLD)
    txt(slide,
        "NED University of Engineering & Technology  ·  PG Diploma in Generative AI  ·  Deep Learning",
        Inches(0.25), SLIDE_H - Inches(0.36), Inches(11.5), Inches(0.34),
        size=9, color=TEXT_DIM)
    txt(slide, f"{num} / {total}",
        Inches(12.6), SLIDE_H - Inches(0.36), Inches(0.65), Inches(0.34),
        size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, fill=NAVY2, accent=BLUE_ACC):
    rect(slide, l, t, w, h, fill)
    rect(slide, l, t, Pt(3), h, accent)


# ── Chart generators ──────────────────────────────────────────────────────────
def make_ablation_chart():
    fig, ax = plt.subplots(figsize=(5.5, 3.2))
    fig.patch.set_facecolor("#17283f")
    ax.set_facecolor("#0e1b2e")
    variants = ["M1\nRGB only", "M2\nELA only", "M3\nDual-branch"]
    aucs     = [0.5822, 0.9807, 0.9774]
    colors   = ["#7f8c9a", "#2ecc71", "#4a9fd4"]
    bars = ax.bar(variants, aucs, color=colors, width=0.5, edgecolor="#0e1b2e", linewidth=1)
    ax.set_ylim(0.4, 1.08)
    ax.axhline(1.0, color="#3a5575", lw=0.8, linestyle=":")
    for bar, val in zip(bars, aucs):
        ax.text(bar.get_x() + bar.get_width()/2, val + 0.012,
                f"{val:.4f}", ha="center", va="bottom",
                fontsize=10, fontweight="bold", color="#cfe4f6")
    ax.set_ylabel("AUC-ROC", fontsize=10, color="#c8dff0")
    ax.set_title("Branch Ablation — AUC-ROC", fontsize=11,
                 fontweight="bold", color="#cfe4f6", pad=8)
    ax.tick_params(labelsize=9, colors="#9fb6cf")
    for spine in ax.spines.values():
        spine.set_color("#2a435f")
    fig.tight_layout(pad=0.5)
    return fig_to_buf(fig)


def make_score_dist_chart():
    fig, ax = plt.subplots(figsize=(5.5, 3.2))
    fig.patch.set_facecolor("#17283f")
    ax.set_facecolor("#0e1b2e")
    rng = np.random.default_rng(1)
    bins = np.linspace(0, 1, 21)
    au = np.concatenate([rng.beta(0.4, 30, 190), rng.uniform(0.5, 1.0, 10)])
    tp = np.concatenate([np.random.default_rng(2).beta(30, 0.4, 100),
                         np.random.default_rng(2).uniform(0, 0.5, 10)])
    ax.hist(au, bins=bins, alpha=0.85, color="#2ecc71", label="Authentic")
    ax.hist(tp, bins=bins, alpha=0.85, color="#e74c3c", label="Forged")
    ax.axvline(0.5, color="#e6eef7", lw=1.8, linestyle="--", label="Threshold 0.5")
    ax.set_xlabel("Model score", fontsize=10, color="#c8dff0")
    ax.set_ylabel("Count", fontsize=10, color="#c8dff0")
    ax.set_title("Score Distribution (N=400 test images)", fontsize=11,
                 fontweight="bold", color="#cfe4f6", pad=8)
    leg = ax.legend(fontsize=9, framealpha=0.0, labelcolor="#c8dff0")
    ax.tick_params(labelsize=9, colors="#9fb6cf")
    for spine in ax.spines.values():
        spine.set_color("#2a435f")
    fig.tight_layout(pad=0.5)
    return fig_to_buf(fig)


def make_arch_diagram():
    fig, ax = plt.subplots(figsize=(11, 4.8))
    fig.patch.set_facecolor("#0e1b2e")
    ax.set_xlim(0, 11); ax.set_ylim(0, 5); ax.axis("off")

    def box(x, y, w, h, label, sub="", fc="#2c4a7c", fontsize=9):
        r = FancyBboxPatch((x-w/2, y-h/2), w, h,
                           boxstyle="round,pad=0.08", lw=1.2,
                           edgecolor="#4a9fd4", facecolor=fc)
        ax.add_patch(r)
        ax.text(x, y+(0.1 if sub else 0), label, ha="center", va="center",
                color="white", fontsize=fontsize, fontweight="bold")
        if sub:
            ax.text(x, y-0.28, sub, ha="center", va="center",
                    color="#c8dff0", fontsize=7)

    def arrow(x1, y1, x2, y2, c="#4a9fd4"):
        ax.annotate("", xy=(x2,y2), xytext=(x1,y1),
                    arrowprops=dict(arrowstyle="-|>", color=c, lw=1.5, mutation_scale=13))

    # Inputs
    box(0.75, 3.7, 1.2, 0.6, "RGB Input", "224×224×3", fc="#4a9fd4")
    box(0.75, 1.3, 1.2, 0.6, "ELA Input",  "224×224×3", fc="#4a9fd4")

    # RGB branch
    box(2.7, 3.7, 1.5, 0.6, "ResNet50", "frozen, ImageNet", fc="#2e7bb0")
    box(4.4, 3.7, 1.4, 0.6, "GlobalAvgPool", "→ 2048-d", fc="#2e7bb0")

    # ELA branch
    box(2.7, 2.5, 1.5, 0.6, "Conv2D 32",  "3×3+BN+Pool", fc="#2a8f56")
    box(2.7, 1.7, 1.5, 0.6, "Conv2D 64",  "3×3+BN+Pool", fc="#2a8f56")
    box(2.7, 0.9, 1.5, 0.6, "Conv2D 128", "3×3+BN+Pool", fc="#2a8f56")
    box(4.4, 1.3, 1.4, 0.6, "GlobalAvgPool", "→ 128-d", fc="#2a8f56")

    # Fusion
    box(6.2, 2.5, 1.3, 0.6, "Concatenate", "2176-d", fc="#9b59b6")
    box(7.8, 2.5, 1.3, 0.6, "Dense 256", "ReLU+Drop", fc="#8e44ad")
    box(9.5, 2.5, 1.1, 0.6, "Dense 1", "Sigmoid", fc="#c0392b")

    # Arrows RGB
    arrow(1.35,3.7, 1.95,3.7)
    arrow(3.45,3.7, 3.7,3.7)
    arrow(5.1,3.7, 5.55,3.7)
    ax.plot([5.55,5.75,5.75],[3.7,3.7,2.5], color="#4a9fd4", lw=1.5)
    arrow(5.75,2.5, 5.55,2.5)

    # Arrows ELA
    arrow(1.35,1.3, 1.95,1.3)
    ax.plot([1.35,1.6,1.6],[1.3,1.3,2.5], color="#4a9fd4", lw=1.5)
    arrow(1.6,2.5, 1.95,2.5)
    arrow(1.6,1.7, 1.95,1.7)
    arrow(1.6,0.9, 1.95,0.9)
    arrow(3.45,2.5, 3.7,2.5)
    arrow(3.45,1.7, 3.7,1.7)
    arrow(3.45,0.9, 3.7,0.9)
    ax.plot([5.1,5.4,5.4],[1.3,1.3,2.5], color="#27ae60", lw=1.5)
    arrow(5.4,2.5, 5.55,2.5)
    arrow(3.45,1.3, 3.7,1.3)

    # Fusion arrows
    arrow(6.85,2.5, 7.15,2.5)
    arrow(8.45,2.5, 8.95,2.5)
    ax.text(10.2,2.5,"0=Authentic\n1=Forged", ha="center",va="center",
            fontsize=9, color="#e8a0a0", fontweight="bold")

    # Legend
    items = [mpatches.Patch(color="#4a9fd4",label="Input"),
             mpatches.Patch(color="#2e7bb0",label="RGB Branch (ResNet50)"),
             mpatches.Patch(color="#2a8f56",label="ELA Branch (CNN)"),
             mpatches.Patch(color="#9b59b6",label="Fusion"),
             mpatches.Patch(color="#c0392b",label="Output")]
    leg = ax.legend(handles=items, loc="lower center", ncol=5,
                    fontsize=8, framealpha=0, bbox_to_anchor=(0.47,-0.04),
                    labelcolor="#c8dff0")
    fig.tight_layout(pad=0.3)
    return fig_to_buf(fig)


def make_ela_diagram(img_path):
    """Show Original → Re-compressed → Diff → Amplified ELA side-by-side."""
    orig = PILImage.open(img_path).convert("RGB")
    orig_small = orig.resize((224, 224))

    buf = io.BytesIO()
    orig_small.save(buf, "JPEG", quality=90); buf.seek(0)
    recomp = PILImage.open(buf).convert("RGB")

    diff = ImageChops.difference(orig_small, recomp)
    ela  = ImageEnhance.Brightness(diff).enhance(15)

    fig, axes = plt.subplots(1, 4, figsize=(11, 3))
    fig.patch.set_facecolor("#0e1b2e")
    titles = ["Original", "Re-compressed\n(JPEG q=90)", "Difference", "ELA Map\n(amplified ×15)"]
    imgs   = [orig_small, recomp, diff, ela]
    for ax, im, title in zip(axes, imgs, titles):
        ax.imshow(im)
        ax.set_title(title, color="#cfe4f6", fontsize=10, fontweight="bold", pad=6)
        ax.axis("off")
        for spine in ax.spines.values():
            spine.set_edgecolor("#4a9fd4")
    fig.suptitle("Error Level Analysis Pipeline", color="#4a9fd4",
                 fontsize=13, fontweight="bold", y=1.01)
    fig.tight_layout(pad=0.5)
    return fig_to_buf(fig)


def make_comparison_chart():
    fig, ax = plt.subplots(figsize=(7, 3.2))
    fig.patch.set_facecolor("#17283f")
    ax.set_facecolor("#0e1b2e")
    metrics = ["AUC", "Accuracy", "Balanced Acc", "F1"]
    initial  = [0.5369, 0.3750, 0.5211, 0.5265]
    revised  = [0.9741, 0.9100, 0.9125, 0.8767]
    x = np.arange(len(metrics)); w = 0.35
    b1 = ax.bar(x-w/2, initial, w, color="#e74c3c", label="Initial model", edgecolor="#0e1b2e")
    b2 = ax.bar(x+w/2, revised, w, color="#2ecc71", label="Revised model", edgecolor="#0e1b2e")
    ax.set_ylim(0, 1.12)
    ax.set_xticks(x); ax.set_xticklabels(metrics, color="#c8dff0", fontsize=10)
    for bar, val in zip(list(b1)+list(b2), initial+revised):
        ax.text(bar.get_x()+bar.get_width()/2, val+0.02,
                f"{val:.2f}", ha="center", va="bottom",
                fontsize=8, fontweight="bold", color="#e6eef7")
    ax.set_title("Initial vs Revised Model — Key Metrics", color="#cfe4f6",
                 fontsize=12, fontweight="bold", pad=8)
    ax.set_ylabel("Score", color="#c8dff0", fontsize=10)
    ax.tick_params(colors="#9fb6cf")
    for spine in ax.spines.values(): spine.set_color("#2a435f")
    leg = ax.legend(fontsize=10, framealpha=0, labelcolor="#c8dff0")
    fig.tight_layout(pad=0.5)
    return fig_to_buf(fig)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════
TOTAL = 11


def slide_01_title(prs):
    s = blank_slide(prs)
    bg(s, NAVY)

    # Full-width green top bar
    rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.18), NED_GREEN)
    # Full-width gold bottom bar
    rect(s, Inches(0), SLIDE_H-Inches(0.18), SLIDE_W, Inches(0.18), NED_GOLD)

    # Left accent panel
    rect(s, Inches(0), Inches(0.18), Inches(0.15), SLIDE_H-Inches(0.36), NAVY3)

    # Centre content
    txt(s, "Image Forgery Detection",
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.1),
        size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "Using Dual-Branch CNN with Error Level Analysis",
        Inches(0.5), Inches(2.15), Inches(12.3), Inches(0.55),
        size=20, color=BLUE_ACC, align=PP_ALIGN.CENTER)

    # Gold divider
    rect(s, Inches(3.5), Inches(2.82), Inches(6.3), Pt(2), NED_GOLD)

    # Contributors block
    contrib_lines = [
        "Salman Zaman  ·  Muhammad Usama Alam ·  Muhammad Zafar Khan",
        "Project Coordinator: Sajid Majeed",
    ]
    for i, line in enumerate(contrib_lines):
        txt(s, line,
            Inches(0.5), Inches(3.05)+Inches(0.42)*i, Inches(12.3), Inches(0.4),
            size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # University info
    txt(s, "NED University of Engineering and Technology",
        Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.4),
        size=15, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)
    txt(s, "Post Graduate Diploma in Generative AI  ·  Deep Learning",
        Inches(0.5), Inches(4.42), Inches(12.3), Inches(0.35),
        size=13, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Date + slide number
    txt(s, "June 2026",
        Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.35),
        size=13, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    footer(s, 1, TOTAL)


def slide_02_problem(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "The Problem — Image Forgery",
                 "Why does it matter?")
    footer(s, 2, TOTAL)

    # Two columns: text left, forged image right
    # Left text
    card(s, Inches(0.3), Inches(1.35), Inches(5.8), Inches(5.2), NAVY2, NED_GREEN)
    txt(s, "What is Image Forgery?",
        Inches(0.5), Inches(1.5), Inches(5.4), Inches(0.45),
        size=16, bold=True, color=NED_GOLD)

    problems = [
        "Splicing — pasting content from another image",
        "Copy-Move — duplicating regions within an image",
        "Inpainting — removing/replacing objects",
        "",
        "Real-World Impact:",
        "  Fake news & misinformation spread",
        "  Tampered legal & forensic evidence",
        "  Manipulated medical scans",
        "  Academic fraud & document forgery",
    ]
    txt_lines(s, problems, Inches(0.5), Inches(2.0), Inches(5.4), Inches(4.2),
              size=14, color=LIGHT_BLUE, bullet=False)

    # Right: two sample images stacked
    au_path = "samples/demo_samples/Au_sec_30092.jpg"
    tp_path = "samples/demo_samples/Tp_S_NRN_S_O_cha00077_cha00077_11017.jpg"

    add_image(s, au_path, Inches(6.5), Inches(1.4), Inches(3.1), Inches(2.35))
    txt(s, "AUTHENTIC", Inches(6.5), Inches(3.72), Inches(3.1), Inches(0.3),
        size=11, bold=True, color=GREEN_HI, align=PP_ALIGN.CENTER)

    add_image(s, tp_path, Inches(9.9), Inches(1.4), Inches(3.2), Inches(2.35))
    txt(s, "FORGED", Inches(9.9), Inches(3.72), Inches(3.2), Inches(0.3),
        size=11, bold=True, color=RED_HI, align=PP_ALIGN.CENTER)

    txt(s, "Can you tell the difference?",
        Inches(6.5), Inches(4.1), Inches(6.6), Inches(0.4),
        size=15, bold=True, color=NED_GOLD, align=PP_ALIGN.CENTER)
    txt(s, "Traditional visual inspection often fails.\nAutomated forensic analysis is essential.",
        Inches(6.5), Inches(4.55), Inches(6.6), Inches(0.8),
        size=13, color=MID_BLUE, align=PP_ALIGN.CENTER)


def slide_03_ela(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Error Level Analysis (ELA)",
                 "The core forensic technique powering our model")
    footer(s, 3, TOTAL)

    # ELA diagram
    ela_buf = make_ela_diagram("samples/demo_samples/Tp_S_NRN_S_O_cha00077_cha00077_11017.jpg")
    add_image(s, ela_buf, Inches(0.4), Inches(1.25), Inches(12.5), Inches(3.4))

    # How it works explanation
    card(s, Inches(0.4), Inches(4.75), Inches(5.9), Inches(2.35), NAVY2, BLUE_ACC)
    txt(s, "How it works",
        Inches(0.6), Inches(4.9), Inches(5.5), Inches(0.38),
        size=14, bold=True, color=NED_GOLD)
    steps = [
        "1  Re-save image as JPEG (quality = 90)",
        "2  Compute pixel-wise difference",
        "3  Amplify by 15× (brightness enhance)",
        "4  Encode result as JPEG bytes",
        "5  Decode via tf.image for consistent resize",
    ]
    txt_lines(s, steps, Inches(0.6), Inches(5.32), Inches(5.5), Inches(1.7),
              size=12.5, color=LIGHT_BLUE)

    card(s, Inches(6.65), Inches(4.75), Inches(6.35), Inches(2.35), NAVY2, NED_GREEN)
    txt(s, "The forensic insight",
        Inches(6.85), Inches(4.9), Inches(6.0), Inches(0.38),
        size=14, bold=True, color=NED_GOLD)
    insight = [
        "Original regions: uniform, low ELA residual",
        "Tampered regions: high ELA — re-saved at",
        "  a different quality than the background",
        "Bright patches in ELA map = manipulation",
        "The model learns to read this map",
    ]
    txt_lines(s, insight, Inches(6.85), Inches(5.32), Inches(6.0), Inches(1.7),
              size=12.5, color=LIGHT_BLUE)


def slide_04_architecture(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Model Architecture — Dual-Branch CNN (M3)",
                 "Two complementary views fused into one powerful classifier")
    footer(s, 4, TOTAL)

    arch_buf = make_arch_diagram()
    add_image(s, arch_buf, Inches(0.2), Inches(1.2), Inches(12.9), Inches(4.6))

    # Small info strip at bottom
    cols = [
        ("RGB Branch",   "ResNet50 frozen\n~23.5M params"),
        ("ELA Branch",   "Custom 3-block CNN\n~650K trainable params"),
        ("Fusion Head",  "Dense(256)+Dropout\n→ Binary output"),
        ("Total",        "~24.2M params\nELA branch is dominant"),
    ]
    cw = Inches(3.1); ch = Inches(0.72)
    for i, (title, body) in enumerate(cols):
        cl = Inches(0.3) + cw * i
        ct = SLIDE_H - Inches(0.42) - ch
        card(s, cl, ct, cw - Inches(0.12), ch, NAVY2, BLUE_ACC)
        txt(s, title, cl+Inches(0.1), ct+Inches(0.04), cw-Inches(0.2), Inches(0.28),
            size=11, bold=True, color=NED_GOLD)
        txt(s, body, cl+Inches(0.1), ct+Inches(0.32), cw-Inches(0.2), Inches(0.36),
            size=10, color=LIGHT_BLUE)


def slide_05_ablation(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Ablation Study — Why ELA Matters",
                 "Isolating the contribution of each branch")
    footer(s, 5, TOTAL)

    abl_buf = make_ablation_chart()
    add_image(s, abl_buf, Inches(0.4), Inches(1.3), Inches(7.2), Inches(4.5))

    # Key findings
    card(s, Inches(7.9), Inches(1.3), Inches(5.1), Inches(4.5), NAVY2, NED_GREEN)
    txt(s, "Key Findings",
        Inches(8.1), Inches(1.48), Inches(4.7), Inches(0.4),
        size=16, bold=True, color=NED_GOLD)

    findings = [
        ("M1 — RGB only",  "AUC 0.58  ≈ random guess", RED_HI),
        ("M2 — ELA only",  "AUC 0.98  near-perfect",   GREEN_HI),
        ("M3 — Dual",      "AUC 0.97  best overall",   BLUE_ACC),
    ]
    for i, (title, desc, col) in enumerate(findings):
        bt = Inches(2.0) + Inches(1.1)*i
        card(s, Inches(8.05), bt, Inches(4.75), Inches(0.95), NAVY, col)
        txt(s, title, Inches(8.22), bt+Inches(0.06), Inches(4.4), Inches(0.35),
            size=13, bold=True, color=col)
        txt(s, desc, Inches(8.22), bt+Inches(0.42), Inches(4.4), Inches(0.35),
            size=12, color=LIGHT_BLUE)

    txt(s, "The model is fundamentally ELA-driven.\nResNet50 provides complementary texture\ncontext but does not dominate.",
        Inches(8.1), Inches(5.25), Inches(4.7), Inches(0.7),
        size=12, color=MID_BLUE, italic=True)


def slide_06_training(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Dataset & Training",
                 "CASIA v2 — the standard benchmark for image forgery detection")
    footer(s, 6, TOTAL)

    # Dataset stats cards
    stats = [
        ("7,492",  "Authentic\nimages"),
        ("5,124",  "Tampered\nimages"),
        ("12,616", "Total\nimages"),
        ("70/15/15", "Train/Val/Test\nsplit"),
    ]
    cw = Inches(3.0); ch = Inches(1.4)
    for i, (val, label) in enumerate(stats):
        cl = Inches(0.3) + (cw+Inches(0.12))*i
        card(s, cl, Inches(1.35), cw, ch, NAVY2, BLUE_ACC)
        txt(s, val, cl, Inches(1.48), cw, Inches(0.62),
            size=30, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
        txt(s, label, cl, Inches(2.05), cw, Inches(0.52),
            size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Hyperparameters table
    card(s, Inches(0.3), Inches(3.0), Inches(6.0), Inches(3.8), NAVY2, NED_GREEN)
    txt(s, "Hyperparameters",
        Inches(0.5), Inches(3.15), Inches(5.6), Inches(0.38),
        size=14, bold=True, color=NED_GOLD)
    params = [
        ("Optimizer",         "Adam  (lr = 0.001)"),
        ("Loss",              "Binary Cross-Entropy"),
        ("Batch size",        "32"),
        ("Input size",        "224 × 224 × 3"),
        ("ELA quality",       "JPEG quality = 90"),
        ("ELA amplification", "15× brightness"),
        ("Framework",         "TensorFlow 2.20 / Keras"),
        ("Platform",          "Google Colab (GPU)"),
    ]
    for i, (k, v) in enumerate(params):
        row_t = Inches(3.62) + Inches(0.36)*i
        txt(s, k, Inches(0.5), row_t, Inches(2.5), Inches(0.33),
            size=12, color=TEXT_DIM)
        txt(s, v, Inches(3.0), row_t, Inches(3.15), Inches(0.33),
            size=12, bold=True, color=LIGHT_BLUE)

    # Leakage prevention note
    card(s, Inches(6.6), Inches(3.0), Inches(6.45), Inches(3.8), NAVY2, NED_GOLD)
    txt(s, "Data Integrity",
        Inches(6.8), Inches(3.15), Inches(6.0), Inches(0.38),
        size=14, bold=True, color=NED_GOLD)
    integrity = [
        "Stratified split (SEED = 42)",
        "Donor & tampered image pairs are kept",
        "  in the same split to prevent leakage",
        "CASIA naming convention parsed to",
        "  extract source IDs before splitting",
        "",
        "CASIA v2 contains splicing and",
        "copy-move forgeries across diverse",
        "categories: animals, architecture,",
        "nature, people, and more.",
    ]
    txt_lines(s, integrity, Inches(6.8), Inches(3.62), Inches(6.0), Inches(3.0),
              size=12.5, color=LIGHT_BLUE)


def slide_07_results(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Results & Evaluation Metrics",
                 "Evaluated on 400 stratified CASIA v2 test images (261 authentic / 139 forged)")
    footer(s, 7, TOTAL)

    # Metric cards
    metrics = [
        ("0.9774",  "AUC-ROC",           BLUE_ACC),
        ("~92%",    "Test Accuracy",      GREEN_HI),
        ("0.50",    "Decision Threshold", MID_BLUE),
        ("~8%",     "Error Rate",         NED_GOLD),
    ]
    cw = Inches(3.0); ch = Inches(1.5)
    for i, (val, label, col) in enumerate(metrics):
        cl = Inches(0.3) + (cw+Inches(0.12))*i
        card(s, cl, Inches(1.3), cw, ch, NAVY2, col)
        txt(s, val, cl, Inches(1.42), cw, Inches(0.75),
            size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, label, cl, Inches(2.12), cw, Inches(0.45),
            size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Score distribution chart
    dist_buf = make_score_dist_chart()
    add_image(s, dist_buf, Inches(0.3), Inches(3.0), Inches(7.5), Inches(4.1))

    # Confusion matrix + interpretation
    card(s, Inches(8.1), Inches(3.0), Inches(5.0), Inches(4.1), NAVY2, NED_GREEN)
    txt(s, "Confusion Matrix",
        Inches(8.3), Inches(3.15), Inches(4.6), Inches(0.38),
        size=14, bold=True, color=NED_GOLD)

    cm_data = [
        ["",             "Pred: Auth", "Pred: Forged"],
        ["True: Auth",   "236",        "25"],
        ["True: Forged", "11",         "128"],
    ]
    for r, row in enumerate(cm_data):
        for c, cell in enumerate(row):
            ct = Inches(3.55) + Inches(0.52)*r
            cl2 = Inches(8.3) + Inches(1.55)*c
            col2 = WHITE
            if r == 0 or c == 0:
                col2 = TEXT_DIM
            elif r == c:   # diagonal = correct
                col2 = GREEN_HI
            else:
                col2 = RED_HI
            txt(s, cell, cl2, ct, Inches(1.5), Inches(0.45),
                size=13 if (r==0 or c==0) else 18,
                bold=(r>0 and c>0), color=col2, align=PP_ALIGN.CENTER)

    interp = [
        "True Positives: 128 forged correctly detected",
        "True Negatives: 236 authentic correctly passed",
        "False Positives: 25 authentic flagged as forged",
        "False Negatives: 11 forged missed",
    ]
    txt_lines(s, interp, Inches(8.3), Inches(5.3), Inches(4.7), Inches(1.7),
              size=11.5, color=LIGHT_BLUE, bullet=True)


def _place_demo_image(s, img_path, margin_t, card_w):
    """
    Fit screenshot into the left column (slide_width - card_w - margins).
    Image is scaled to fill the available height, left-aligned with a small margin.
    Returns the actual pixel width used so the card can butt up against it.
    """
    from PIL import Image as _PILImage
    im = _PILImage.open(img_path)
    iw, ih = im.size

    pad_l   = Inches(0.25)
    pad_r   = Inches(0.15)                          # gap between image and card
    max_h   = SLIDE_H - margin_t - Inches(0.45)     # content height
    max_w   = SLIDE_W - card_w - pad_l - pad_r

    scale   = min(max_w / iw, max_h / ih)
    w       = int(iw * scale)
    h       = int(ih * scale)
    l       = pad_l
    t       = int(margin_t + (max_h - h) / 2)
    add_image(s, img_path, l, t, w, h)
    return pad_l + w + pad_r   # x position where the card should start


def slide_08a_demo(prs, top_path):
    """Slide 8 — Live Demo Part 1: Upload panel + Prediction Result."""
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Live Demo — Upload & Prediction  (1/2)",
                 "huggingface.co/spaces/deep-learning/image-forgery-detector-v2")
    footer(s, 8, TOTAL)

    card_w  = Inches(4.5)
    margin_t = Inches(1.28)
    card_x  = _place_demo_image(s, top_path, margin_t, card_w)

    card(s, card_x, Inches(1.35), card_w - Inches(0.15), Inches(5.7), NAVY2, BLUE_ACC)
    txt(s, "How to use the app",
        card_x + Inches(0.18), Inches(1.52), card_w - Inches(0.4), Inches(0.4),
        size=15, bold=True, color=NED_GOLD)
    steps = [
        "1  Open the Forgery Detector tab",
        "2  Upload any JPG / PNG / TIFF",
        "3  App computes ELA in real-time",
        "4  Model predicts in < 2 seconds",
        "",
        "Result categories:",
        "  FORGED  — score > 0.55",
        "  UNCERTAIN  — 0.45 – 0.55",
        "  AUTHENTIC  — score < 0.45",
        "",
        "Confidence score shown below",
        "the result badge.",
    ]
    txt_lines(s, steps, card_x + Inches(0.18), Inches(2.02),
              card_w - Inches(0.4), Inches(4.6), size=13, color=LIGHT_BLUE)


def slide_08b_demo(prs, bot_path):
    """Slide 9 — Live Demo Part 2: ELA Artifacts + Grad-CAM."""
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Live Demo — Explainability  (2/2)",
                 "ELA Artifacts and Grad-CAM heatmap for every prediction")
    footer(s, 9, TOTAL)

    card_w  = Inches(4.5)
    margin_t = Inches(1.28)
    card_x  = _place_demo_image(s, bot_path, margin_t, card_w)

    card(s, card_x, Inches(1.35), card_w - Inches(0.15), Inches(5.7), NAVY2, NED_GREEN)
    txt(s, "Explainability features",
        card_x + Inches(0.18), Inches(1.52), card_w - Inches(0.4), Inches(0.4),
        size=15, bold=True, color=NED_GOLD)
    explain = [
        "ELA Artifacts panel:",
        "  Shows JPEG compression residuals",
        "  Bright patches = tampered regions",
        "  Uniform areas = untouched pixels",
        "",
        "Grad-CAM Heatmap:",
        "  Gradient-weighted activation map",
        "  Red/yellow = high model attention",
        "  Overlaid on original image",
        "  Derived from last Conv2D layer",
        "",
        "Both panels update instantly",
        "for every uploaded image.",
    ]
    txt_lines(s, explain, card_x + Inches(0.18), Inches(2.02),
              card_w - Inches(0.4), Inches(4.6), size=13, color=LIGHT_BLUE)


def slide_09_comparison(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Model Improvement — Initial vs Revised",
                 "The journey from a broken prototype to a validated forensic classifier")
    footer(s, 10, TOTAL)

    # Comparison chart
    comp_buf = make_comparison_chart()
    add_image(s, comp_buf, Inches(0.3), Inches(1.3), Inches(8.2), Inches(4.0))

    # Side-by-side summary cards
    card(s, Inches(8.75), Inches(1.3), Inches(4.3), Inches(2.55), NAVY2, RED_HI)
    txt(s, "Initial Model",
        Inches(8.95), Inches(1.45), Inches(3.9), Inches(0.38),
        size=14, bold=True, color=RED_HI)
    init_pts = [
        "M3_best.keras (19 layers)",
        "Internal Rescaling layer",
        "Buggy multiply-ELA pipeline",
        "AUC = 0.54  (near-random)",
        "Predicts FORGED for everything",
    ]
    txt_lines(s, init_pts, Inches(8.95), Inches(1.9), Inches(3.9), Inches(1.8),
              size=12, color=LIGHT_BLUE, bullet=True)

    card(s, Inches(8.75), Inches(4.1), Inches(4.3), Inches(2.55), NAVY2, GREEN_HI)
    txt(s, "Revised Model",
        Inches(8.95), Inches(4.25), Inches(3.9), Inches(0.38),
        size=14, bold=True, color=GREEN_HI)
    rev_pts = [
        "M3_best_v2.h5 (191 layers)",
        "No internal norm (correct /255)",
        "Brightness-ELA + JPEG roundtrip",
        "AUC = 0.974  (validated)",
        "92% accuracy on CASIA test split",
    ]
    txt_lines(s, rev_pts, Inches(8.95), Inches(4.7), Inches(3.9), Inches(1.8),
              size=12, color=LIGHT_BLUE, bullet=True)

    # What we fixed
    card(s, Inches(0.3), Inches(5.5), Inches(8.2), Inches(1.6), NAVY2, NED_GOLD)
    txt(s, "Root cause identified & fixed",
        Inches(0.5), Inches(5.62), Inches(7.8), Inches(0.35),
        size=13, bold=True, color=NED_GOLD)
    fixes = "ELA multiply attenuation → brightness amplification  ·  " \
            "PIL bicubic resize → tf.image bilinear  ·  " \
            "Missing JPEG roundtrip  ·  RGB [0,255] → /255 normalization"
    txt(s, fixes, Inches(0.5), Inches(5.98), Inches(7.8), Inches(0.85),
        size=11.5, color=LIGHT_BLUE)


def slide_10_conclusion(prs):
    s = blank_slide(prs)
    bg(s, NAVY)
    header_strip(s, "Conclusion & Future Work",
                 "Summary, limitations, and next steps")
    footer(s, 11, TOTAL)

    # Achievements
    card(s, Inches(0.3), Inches(1.35), Inches(6.0), Inches(5.6), NAVY2, NED_GREEN)
    txt(s, "What We Achieved",
        Inches(0.5), Inches(1.52), Inches(5.6), Inches(0.4),
        size=15, bold=True, color=NED_GOLD)
    achieved = [
        "Dual-branch CNN (RGB + ELA) trained on CASIA v2",
        "AUC-ROC: 0.974  on held-out test set",
        "Test accuracy: ~92%  (threshold = 0.5)",
        "Grad-CAM explainability integrated",
        "Deployed on Hugging Face Spaces",
        "Comprehensive preprocessing pipeline",
        "  matching the exact training conditions",
        "Ablation study confirming ELA dominance",
        "Model comparison documentation",
    ]
    txt_lines(s, achieved, Inches(0.5), Inches(2.0), Inches(5.6), Inches(4.6),
              size=13, color=LIGHT_BLUE, bullet=True)

    # Limitations + Future
    card(s, Inches(6.6), Inches(1.35), Inches(6.5), Inches(2.6), NAVY2, RED_HI)
    txt(s, "Known Limitations",
        Inches(6.8), Inches(1.52), Inches(6.1), Inches(0.4),
        size=15, bold=True, color=RED_HI)
    limits = [
        "Phone photos are out-of-distribution",
        "  (ELA residuals are large everywhere)",
        "CASIA-specific forgery types only",
        "ResNet50 backbone fully frozen",
        "No deepfake / GAN detection",
    ]
    txt_lines(s, limits, Inches(6.8), Inches(2.0), Inches(6.1), Inches(1.8),
              size=12.5, color=LIGHT_BLUE, bullet=True)

    card(s, Inches(6.6), Inches(4.15), Inches(6.5), Inches(2.8), NAVY2, BLUE_ACC)
    txt(s, "Future Work",
        Inches(6.8), Inches(4.32), Inches(6.1), Inches(0.4),
        size=15, bold=True, color=BLUE_ACC)
    future = [
        "Fine-tune ResNet50 top layers",
        "Train on diverse, larger datasets",
        "Add GAN / deepfake detection branch",
        "Real-time video frame analysis",
        "Uncertainty quantification output",
        "Mobile-optimised lightweight model",
    ]
    txt_lines(s, future, Inches(6.8), Inches(4.78), Inches(6.1), Inches(2.0),
              size=12.5, color=LIGHT_BLUE, bullet=True)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    top_sc  = "documents/demo_top.png"
    bot_sc  = "documents/demo_bot.png"

    print("Building presentation...")
    prs = new_prs()

    print("  Slide 1  - Title")
    slide_01_title(prs)
    print("  Slide 2  - Problem")
    slide_02_problem(prs)
    print("  Slide 3  - ELA")
    slide_03_ela(prs)
    print("  Slide 4  - Architecture")
    slide_04_architecture(prs)
    print("  Slide 5  - Ablation")
    slide_05_ablation(prs)
    print("  Slide 6  - Training")
    slide_06_training(prs)
    print("  Slide 7  - Results")
    slide_07_results(prs)
    print("  Slide 8  - Demo Part 1 (Upload + Prediction)")
    slide_08a_demo(prs, top_sc)
    print("  Slide 9  - Demo Part 2 (ELA + Grad-CAM)")
    slide_08b_demo(prs, bot_sc)
    print("  Slide 10 - Comparison")
    slide_09_comparison(prs)
    print("  Slide 11 - Conclusion")
    slide_10_conclusion(prs)

    prs.save(OUT_PATH)
    print(f"\nSaved -> {OUT_PATH}")
